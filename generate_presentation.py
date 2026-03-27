"""
Generate thesis presentation: RCuckoo vs OfflineState.
Systems-conference style (USENIX/OSDI inspired).
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BASEDIR = os.path.dirname(os.path.abspath(__file__))

# ── Palette ───────────────────────────────────────────────────
CHARCOAL    = RGBColor(0x2D, 0x2D, 0x2D)
WARM_WHITE  = RGBColor(0xFA, 0xFA, 0xFA)
STEEL       = RGBColor(0x37, 0x47, 0x4F)
TEAL        = RGBColor(0x00, 0x83, 0x8F)
MUTED_RED   = RGBColor(0xD3, 0x2F, 0x2F)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY   = RGBColor(0x33, 0x33, 0x33)
MED_GRAY    = RGBColor(0x77, 0x77, 0x77)
LIGHT_GRAY  = RGBColor(0xBB, 0xBB, 0xBB)
GREEN       = RGBColor(0x2E, 0x7D, 0x32)
ORANGE      = RGBColor(0xEF, 0x6C, 0x00)
TEAL_LIGHT  = RGBColor(0xE0, 0xF2, 0xF1)
STEEL_LIGHT = RGBColor(0xEC, 0xEF, 0xF1)
SHADOW_CLR  = RGBColor(0xCC, 0xCC, 0xCC)

FONT_BODY = "Calibri"
FONT_MONO = "Consolas"
FOOTER_TEXT = "Ilan K. \u2014 Open University of Israel"


# ── Helpers ───────────────────────────────────────────────────

def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _footer(slide, num):
    # thin line
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(0.5), Inches(7.05), Inches(9.0), Pt(0.75))
    ln.fill.solid()
    ln.fill.fore_color.rgb = LIGHT_GRAY
    ln.line.fill.background()
    # left text
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(6), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.text = FOOTER_TEXT
    p.font.size = Pt(9)
    p.font.name = FONT_BODY
    p.font.color.rgb = MED_GRAY
    # right number
    tb2 = slide.shapes.add_textbox(Inches(8.5), Inches(7.1), Inches(1), Inches(0.3))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = str(num)
    p2.font.size = Pt(9)
    p2.font.name = FONT_BODY
    p2.font.color.rgb = MED_GRAY
    p2.alignment = PP_ALIGN.RIGHT


def _header_bar(slide, color=CHARCOAL):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(0), Inches(0), Inches(10), Inches(0.08))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()


def _title(slide, text, color=CHARCOAL, top=Inches(0.25), size=26):
    tb = slide.shapes.add_textbox(Inches(0.6), top, Inches(8.8), Inches(0.7))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.name = FONT_BODY
    p.font.color.rgb = color


def _bullets(slide, items, top=Inches(1.15), left=Inches(0.6), width=Inches(8.8),
             height=Inches(5.3), fs=16, sub=None, fc=DARK_GRAY, dash=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    pre = "\u2013  " if dash else ""
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = pre + item
        p.font.size = Pt(fs)
        p.font.name = FONT_BODY
        p.font.color.rgb = fc
        p.space_after = Pt(7)
        if sub and i in sub:
            for sb in sub[i]:
                sp = tf.add_paragraph()
                sp.text = "      " + sb
                sp.font.size = Pt(fs - 2)
                sp.font.name = FONT_BODY
                sp.font.color.rgb = MED_GRAY
                sp.space_after = Pt(3)


def _note(slide, text, top=Inches(6.6)):
    tb = slide.shapes.add_textbox(Inches(0.6), top, Inches(8.8), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.italic = True
    p.font.name = FONT_BODY
    p.font.color.rgb = MED_GRAY


def _box(slide, left, top, w, h, text, bg, fs=12, bold=False, fc=WHITE,
         border=None, shadow=False, font=None):
    if shadow:
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    left + Pt(3), top + Pt(3), w, h)
        sh.fill.solid()
        sh.fill.fore_color.rgb = SHADOW_CLR
        sh.line.fill.background()
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = bg
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    tf = s.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(fs)
        p.font.bold = bold
        p.font.color.rgb = fc
        p.font.name = font or FONT_BODY
        p.alignment = PP_ALIGN.CENTER
    return s


def _conn(slide, x1, y1, x2, y2, color=MED_GRAY, width=1.5):
    c = slide.shapes.add_connector(1, x1, y1, x2, y2)
    c.line.color.rgb = color
    c.line.width = Pt(width)


def _insight_strip(slide, lines, top=Inches(5.55), bg=TEAL_LIGHT):
    if isinstance(lines, str):
        lines = [lines]
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(0), top, Inches(10), Inches(1.4))
    s.fill.solid()
    s.fill.fore_color.rgb = bg
    s.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.6), top + Inches(0.1),
                                  Inches(8.8), Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "\u2013  " + line
        p.font.size = Pt(12)
        p.font.name = FONT_BODY
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(2)


def _image(slide, filename, top=Inches(0.85), height=Inches(4.55)):
    path = os.path.join(BASEDIR, filename)
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(0.3), top,
                                 width=Inches(9.4), height=height)


def _section_slide(prs, text, subtitle, bg_color):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, bg_color)
    tb = sl.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.name = FONT_BODY
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(18)
    p2.font.name = FONT_BODY
    p2.font.color.rgb = RGBColor(0xCC, 0xDD, 0xDD)
    p2.alignment = PP_ALIGN.LEFT
    p2.space_before = Pt(12)
    return sl


def _content_slide(prs, sn, title_text, title_color=CHARCOAL, bar_color=CHARCOAL,
                    size=26):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header_bar(sl, bar_color)
    _title(sl, title_text, title_color, size=size)
    _footer(sl, sn)
    return sl


# ── Build ─────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    sn = 0

    # ━━ 1. TITLE ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    sn += 1
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, CHARCOAL)

    tb = sl.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Client-Side Caching for\nDisaggregated Key-Value Stores"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.name = FONT_BODY
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    p2 = tf.add_paragraph()
    p2.text = "RCuckoo vs. OfflineState"
    p2.font.size = Pt(22)
    p2.font.name = FONT_BODY
    p2.font.color.rgb = TEAL
    p2.alignment = PP_ALIGN.LEFT
    p2.space_before = Pt(16)

    # accent line
    ln = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0.8), Inches(4.2), Inches(2.5), Inches(0.04))
    ln.fill.solid()
    ln.fill.fore_color.rgb = TEAL
    ln.line.fill.background()

    tb2 = sl.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(8), Inches(1.5))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    for txt, sz, col in [("Ilan K.", 18, WHITE),
                         ("The Open University of Israel", 14, RGBColor(0xAA, 0xAA, 0xAA)),
                         ("M.Sc. Thesis", 13, RGBColor(0x88, 0x88, 0x88))]:
        p = tf2.add_paragraph()
        p.text = txt
        p.font.size = Pt(sz)
        p.font.name = FONT_BODY
        p.font.color.rgb = col
        p.space_before = Pt(4)

    # ━━ 2. OUTLINE ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    sn += 1
    sl = _content_slide(prs, sn, "Outline")
    _bullets(sl, [
        "Motivation: cross-rack disaggregation challenges",
        "Background: RDMA-based key-value stores",
        "RCuckoo design: architecture, operations, configuration",
        "OfflineState design: caching, bloom filters, read/write flows",
        "Evaluation methodology",
        "Results: throughput, latency, RDMA efficiency",
        "Sensitivity analysis",
        "Conclusions & future work",
    ], fs=17, dash=False)

    # ━━ 3. MOTIVATION ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    sn += 1
    sl = _content_slide(prs, sn, "Motivation")
    _bullets(sl, [
        "Disaggregated memory pools network-attached resources",
        "RDMA: low-latency remote access (~1 us same-rack)",
        "Cross-rack latency is 2\u20135x higher than intra-rack",
        "Existing KVS designs assume same-rack co-location",
        "Scaling beyond one rack increases server access cost",
        "Opportunity: exploit peer (1 us) vs server (3 us) asymmetry",
    ], sub={
        2: ["Peer-to-peer (same rack): ~1 us",
            "Client-to-server (cross-rack): ~3 us"],
    })

    # Latency diagram right side
    _box(sl, Inches(6.5), Inches(4.6), Inches(3.0), Inches(1.8),
         "Latency Model\n\nLocal cache: 0 ticks\nPeer (same rack): 1 tick\nServer (cross-rack): 3 ticks",
         WARM_WHITE, 12, False, DARK_GRAY, border=TEAL, shadow=True)

    # ━━ 4. BACKGROUND ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    sn += 1
    sl = _content_slide(prs, sn, "Background: Disaggregated Key-Value Stores")
    _bullets(sl, [
        "Clients access passive remote memory via one-sided RDMA",
        "No server-side CPU in the data path (fully disaggregated)",
        "Key tension: serialization for concurrent writes",
        "Lock-based (RCuckoo, Sherman): inlined values, fine-grained locks",
        "Lock-free (FUSEE, Clover): CAS on 8-byte entries, pointer-based",
        "RCuckoo [Grant & Snoeren, ATC'25]: highest mixed-workload throughput",
    ], sub={
        3: ["Locks enable cuckoo hashing with single-RTT reads for small values"],
        4: ["Extent-based storage requires extra round trips"],
    })
    _note(sl, 'Grant, S. and Snoeren, A. C. "Cuckoo for Clients: Disaggregated Cuckoo Hashing." USENIX ATC 2025.')

    # ━━ 5. SECTION: RCUCKOO ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    sn += 1
    _section_slide(prs, "RCuckoo",
                   "Lock-Based Disaggregated Cuckoo Hashing\n[Grant & Snoeren, USENIX ATC 2025]",
                   STEEL)

    # ━━ 6. RCUCKOO ARCHITECTURE ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    sn += 1
    sl = _content_slide(prs, sn, "RCuckoo: Architecture", STEEL, STEEL)

    # Left column
    items = [
        ("Index Table (Main Memory)",
         "Rows of 8 entries (32b key + 32b value)\n"
         "8-bit version + 64-bit CRC per row\n"
         "RDMA-registered, one-sided read/write"),
        ("Lock Table (NIC Memory, 256KB)",
         "1-bit locks; each covers 16 rows\n"
         "MCAS: up to 64 locks in one atomic op\n"
         "NIC memory: 3x lower CAS latency"),
        ("Dependent Hashing",
         "L1(K) = h1(K) mod T\n"
         "L2(K) = L1 + h2(K) mod floor(f^(f+Z(h3(K)))) mod T\n"
         "f=2.3 bounds cuckoo span; Z(x) = trailing zeros"),
        ("Client Cache (64KB)",
         "Caches rows for speculative BFS search\n"
         "Populated on every read and lock acquire"),
    ]
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(5.0), Inches(5.8))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for t, d in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        if not first:
            p.space_before = Pt(10)
        first = False
        p.text = t
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.name = FONT_BODY
        p.font.color.rgb = STEEL
        for line in d.split('\n'):
            p2 = tf.add_paragraph()
            p2.text = line
            p2.font.size = Pt(11)
            p2.font.name = FONT_BODY
            p2.font.color.rgb = DARK_GRAY

    # Right: stacked diagram
    _box(sl, Inches(5.8), Inches(1.2), Inches(3.8), Inches(0.5),
         "NIC Memory (256KB)", STEEL, 12, True, shadow=True)
    _box(sl, Inches(6.0), Inches(1.8), Inches(3.4), Inches(0.4),
         "Lock Table  [bit vector]", RGBColor(0x54, 0x6E, 0x7A), 10)
    _box(sl, Inches(5.8), Inches(2.6), Inches(3.8), Inches(0.5),
         "Main Memory (Server)", STEEL, 12, True, shadow=True)
    _box(sl, Inches(6.0), Inches(3.2), Inches(3.4), Inches(0.4),
         "Index Table  [rows \u00d7 8 entries]", RGBColor(0x54, 0x6E, 0x7A), 10)
    _box(sl, Inches(6.0), Inches(3.7), Inches(3.4), Inches(0.4),
         "Extent Memory", RGBColor(0x54, 0x6E, 0x7A), 10)
    _box(sl, Inches(5.8), Inches(4.5), Inches(3.8), Inches(0.5),
         "Client", GREEN, 12, True, shadow=True)
    _box(sl, Inches(6.0), Inches(5.1), Inches(1.6), Inches(0.4),
         "64KB Cache", RGBColor(0x3E, 0x8E, 0x41), 10)
    _box(sl, Inches(7.8), Inches(5.1), Inches(1.6), Inches(0.4),
         "Hash Funcs", RGBColor(0x3E, 0x8E, 0x41), 10)
    _conn(sl, Inches(7.7), Inches(4.45), Inches(7.7), Inches(4.18))

    # ━━ 7. RCUCKOO OPERATIONS ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    sn += 1
    sl = _content_slide(prs, sn, "RCuckoo: Operations", STEEL, STEEL)

    # Three horizontal flow strips
    ops = [
        ("READ", "1 RTT", GREEN,
         ["compute L1, L2", "RDMA read rows", "CRC validate", "return value"]),
        ("UPDATE", "2 RTTs", ORANGE,
         ["MCAS lock", "read rows", "write value", "release locks"]),
        ("INSERT", "2+ RTTs", MUTED_RED,
         ["BFS search (local)", "MCAS lock", "verify path", "write + unlock"]),
    ]

    for i, (name, cost, color, steps) in enumerate(ops):
        y = Inches(1.2) + Inches(i * 1.55)
        # Label
        _box(sl, Inches(0.3), y, Inches(1.5), Inches(0.55),
             f"{name}\n{cost}", color, 12, True, shadow=True)
        # Step boxes with arrows
        for j, step in enumerate(steps):
            x = Inches(2.1) + Inches(j * 2.0)
            _box(sl, x, y, Inches(1.7), Inches(0.55),
                 step, RGBColor(0xF5, 0xF5, 0xF5), 10, False, DARK_GRAY,
                 border=color)
            if j < len(steps) - 1:
                _conn(sl, x + Inches(1.75), y + Inches(0.27),
                      x + Inches(1.95), y + Inches(0.27), color)

    # Contention note
    _box(sl, Inches(0.3), Inches(5.85), Inches(9.4), Inches(0.9),
         "Lock contention is the primary bottleneck at scale. "
         "Zipf skew creates hot rows \u2192 MCAS retries grow with client count. "
         "At 320 clients YCSB-A: lock saturation limits throughput to ~20 MOPS.",
         STEEL_LIGHT, 11, False, DARK_GRAY, border=STEEL)

    # ━━ 8. RCUCKOO CONFIG ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    sn += 1
    sl = _content_slide(prs, sn, "RCuckoo: Configuration Choices", STEEL, STEEL)

    rows = [
        ("Parameter", "Value", "Rationale"),
        ("Locality f", "2.3", ">95% fill; 68% of keys within 5 rows (Fig.4)"),
        ("Rows/lock", "16", "Sweet spot: single MCAS for 99% of paths (Fig.5)"),
        ("Locks/MCAS", "64", "ConnectX-5 CAS on 64-bit regions"),
        ("Prepopulate", "90%", "Steady-state conditions, longer cuckoo paths"),
        ("Search depth", "5", "Sufficient for 99%+ insertions at f=2.3"),
        ("Client cache", "64KB", "Speculative BFS; caching whole index has diminishing returns"),
        ("Entry size", "8B", "32-bit key + 32-bit value, inlined"),
    ]

    tb = sl.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9.0), Inches(5.5))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (c1, c2, c3) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{c1:<18s}  {c2:<8s}  {c3}"
        p.font.size = Pt(13) if i > 0 else Pt(14)
        p.font.name = FONT_MONO
        p.font.bold = (i == 0)
        p.font.color.rgb = STEEL if i == 0 else DARK_GRAY
        p.space_after = Pt(6)
        if i == 0:
            p2 = tf.add_paragraph()
            p2.text = "\u2500" * 70
            p2.font.size = Pt(8)
            p2.font.name = FONT_MONO
            p2.font.color.rgb = LIGHT_GRAY

    # ━━ 9. SECTION: OFFLINESTATE ━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    sn += 1
    _section_slide(prs, "OfflineState",
                   "Client-Side Caching with Peer Cooperation",
                   TEAL)

    # ━━ 10. OFFLINESTATE ARCHITECTURE ━━━━━━━━━━━━━━━━━━━━━━━━
    sn += 1
    sl = _content_slide(prs, sn, "OfflineState: Architecture", TEAL, TEAL)

    # Three-tier horizontal
    tiers = [
        (Inches(0.5), "Tier 1: Local\nLRU Cache\n256KB", "0 ticks", GREEN),
        (Inches(3.5), "Tier 2: Peer\nBloom + Broadcast", "1 tick", RGBColor(0x01, 0x57, 0x9B)),
        (Inches(6.5), "Tier 3: Server\nRDMA Read/Write", "3 ticks", MUTED_RED),
    ]
    for x, txt, cost, col in tiers:
        _box(sl, x, Inches(1.2), Inches(2.7), Inches(1.2),
             txt, col, 13, True, shadow=True)
        _box(sl, x + Inches(0.4), Inches(2.5), Inches(1.9), Inches(0.35),
             cost, RGBColor(0xF5, 0xF5, 0xF5), 11, True, DARK_GRAY, border=col)

    _conn(sl, Inches(3.25), Inches(1.8), Inches(3.45), Inches(1.8), MED_GRAY, 2)
    _conn(sl, Inches(6.25), Inches(1.8), Inches(6.45), Inches(1.8), MED_GRAY, 2)

    for x, txt in [(Inches(3.1), "miss"), (Inches(6.1), "miss / FP")]:
        tb = sl.shapes.add_textbox(x, Inches(1.4), Inches(0.8), Inches(0.25))
        p = tb.text_frame.paragraphs[0]
        p.text = txt
        p.font.size = Pt(9)
        p.font.name = FONT_BODY
        p.font.color.rgb = MUTED_RED

    _bullets(sl, [
        "Per-client LRU: 256KB = ~32,000 entries (8B each)",
        "Local bloom: tracks own cache (no false negatives)",
        "Peer bloom: bitwise OR of all local blooms",
        "Offline sync: rebuilds + distributes blooms every 100 ticks",
        "Lock-free writes: direct RDMA, zero contention at any scale",
    ], top=Inches(3.2), fs=15, sub={
        2: ["Auto-scaled: m grows with sqrt(num_clients) to prevent saturation"],
    })

    # ━━ 11. OFFLINESTATE READ FLOW ━━━━━━━━━━━━━━━━━━━━━━━━━━
    sn += 1
    sl = _content_slide(prs, sn, "OfflineState: Read Operation Flow", TEAL, TEAL)

    cx, bw = Inches(3.0), Inches(4.0)
    steps = [
        (Inches(0.85), "read(key)", CHARCOAL),
        (Inches(1.55), "Check local LRU cache", GREEN),
        (Inches(2.35), "Check peer bloom filter", RGBColor(0x01, 0x57, 0x9B)),
        (Inches(3.15), "Broadcast to all peers (1 RTT)", RGBColor(0x01, 0x57, 0x9B)),
        (Inches(3.95), "Server RDMA read (3 ticks)", MUTED_RED),
        (Inches(4.65), "Cache result + update bloom", RGBColor(0x55, 0x55, 0x55)),
    ]
    for y, txt, col in steps:
        _box(sl, cx, y, bw, Inches(0.45), txt, col, 12, True, shadow=True)

    # Arrows between steps
    for i in range(len(steps) - 1):
        y_from = steps[i][0] + Inches(0.48)
        y_to = steps[i + 1][0] - Inches(0.03)
        _conn(sl, Inches(5.0), y_from, Inches(5.0), y_to, MED_GRAY)

    # Side annotations
    annotations = [
        (Inches(1.6), Inches(0.2), "HIT \u2192 return (0 ticks)", GREEN),
        (Inches(2.35), Inches(0.2), "MISS \u2193", MUTED_RED),
        (Inches(2.4), Inches(7.2), "NO \u2192 skip to server", MED_GRAY),
        (Inches(3.15), Inches(0.2), "MAYBE \u2193", ORANGE),
        (Inches(3.2), Inches(7.2), "PEER HIT \u2192 cache + return (1 tick)", RGBColor(0x01, 0x57, 0x9B)),
        (Inches(3.2), Inches(0.2), " ", WHITE),
        (Inches(4.0), Inches(7.2), "FALSE POSITIVE \u2192 server", MUTED_RED),
    ]
    for y, x, txt, col in annotations:
        tb = sl.shapes.add_textbox(x, y, Inches(2.8), Inches(0.3))
        p = tb.text_frame.paragraphs[0]
        p.text = txt
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.name = FONT_BODY
        p.font.color.rgb = col

    _box(sl, Inches(0.3), Inches(5.4), Inches(9.4), Inches(1.2),
         "Read latency (Zipf 0.99, 320 clients):  "
         "~31% local (0t) + ~25% peer (1t) + ~44% server (3t)  \u2192  "
         "avg ~1.6 ticks.  56% of reads bypass server entirely.",
         TEAL_LIGHT, 12, False, DARK_GRAY, border=TEAL)

    # ━━ 12. OFFLINESTATE WRITE FLOW ━━━━━━━━━━━━━━━━━━━━━━━━━
    sn += 1
    sl = _content_slide(prs, sn, "OfflineState: Write Operation Flow", TEAL, TEAL)

    # Simple 3-step flow
    wsteps = [
        (Inches(1.5), "write(key, value)", CHARCOAL),
        (Inches(2.4), "One-sided RDMA write to server (3 ticks)", MUTED_RED),
        (Inches(3.3), "Update local cache if key present", GREEN),
        (Inches(4.2), "Done \u2014 no locks, no coordination", RGBColor(0x55, 0x55, 0x55)),
    ]
    for y, txt, col in wsteps:
        _box(sl, Inches(2.5), y, Inches(5.0), Inches(0.55),
             txt, col, 14, True, shadow=True)
    for i in range(len(wsteps) - 1):
        _conn(sl, Inches(5.0), wsteps[i][0] + Inches(0.58),
              Inches(5.0), wsteps[i + 1][0] - Inches(0.03))

    # Comparison box
    _box(sl, Inches(0.5), Inches(5.3), Inches(9.0), Inches(1.4),
         "Write cost comparison:\n"
         "RCuckoo: MCAS lock + read + write + unlock = 2+ RTTs, contention retries at scale\n"
         "OfflineState: 1 RDMA write = fixed cost, zero contention\n"
         "At 320 clients YCSB-A: RCuckoo ~20 MOPS vs OfflineState ~148 MOPS (7.4\u00d7)",
         STEEL_LIGHT, 12, False, DARK_GRAY, border=STEEL)

    # ━━ 13. BLOOM COOPERATION ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    sn += 1
    sl = _content_slide(prs, sn, "Bloom Filter Peer Cooperation", TEAL, TEAL)

    # Left: protocol
    sections = [
        ("Per-Client State", [
            "Local bloom: keys in own LRU (exact, no FN)",
            "Peer bloom: merged from all clients (read-only, has FP)",
        ]),
        ("Sync Protocol (every 100 ticks)", [
            "1. Rebuild each local bloom from current cache",
            "2. Compute bitwise OR across all N local blooms",
            "3. Distribute merged peer bloom to every client",
        ]),
        ("Auto-Scaling", [
            "Fixed bloom saturates at high N (FPR > 40%)",
            "Scale: n = cache_entries \u00d7 ceil(sqrt(N))",
            "m = \u2013(n \u00b7 ln(p)) / ln(2)\u00b2,  k = (m/n) \u00b7 ln(2)",
            "Result: FPR drops 41% \u2192 3.5% at 320 clients",
        ]),
    ]
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(5.0), Inches(5.5))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for title, items in sections:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        if not first:
            p.space_before = Pt(10)
        first = False
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.name = FONT_BODY
        p.font.color.rgb = TEAL
        for item in items:
            p2 = tf.add_paragraph()
            p2.text = "  \u2013  " + item
            p2.font.size = Pt(11)
            p2.font.name = FONT_BODY
            p2.font.color.rgb = DARK_GRAY

    # Right: diagram
    for j, (y, txt) in enumerate([
        (Inches(1.2), "Client 1\nBloom_1"),
        (Inches(1.2), "Client 2\nBloom_2"),
        (Inches(2.2), "Client 3\nBloom_3"),
        (Inches(2.2), "Client N\nBloom_N"),
    ]):
        x = Inches(5.8) if j % 2 == 0 else Inches(7.8)
        _box(sl, x, y, Inches(1.7), Inches(0.75), txt,
             RGBColor(0x01, 0x57, 0x9B), 10, shadow=True)

    _box(sl, Inches(6.1), Inches(3.3), Inches(3.0), Inches(0.55),
         "Sync Worker", GREEN, 11, True, shadow=True)
    _box(sl, Inches(6.1), Inches(4.15), Inches(3.0), Inches(0.55),
         "Peer = OR(B1 .. BN)", ORANGE, 11, True, shadow=True)
    _box(sl, Inches(6.1), Inches(5.0), Inches(3.0), Inches(0.55),
         "Distribute to all", RGBColor(0x55, 0x55, 0x55), 11, shadow=True)

    _conn(sl, Inches(7.6), Inches(3.0), Inches(7.6), Inches(3.25))
    _conn(sl, Inches(7.6), Inches(3.9), Inches(7.6), Inches(4.1))
    _conn(sl, Inches(7.6), Inches(4.75), Inches(7.6), Inches(4.95))

    # ━━ 14. OFFLINESTATE CONFIG ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    sn += 1
    sl = _content_slide(prs, sn, "OfflineState: Configuration Choices", TEAL, TEAL)

    rows = [
        ("Parameter", "Value", "Rationale"),
        ("Cache size", "256KB", "~32K entries; saturates local hit rate (Plot 9)"),
        ("Sync interval", "100 ticks", "Robust 10\u20131000 range; FPR rises beyond (Plot 8)"),
        ("Bloom FPR target", "1%", "Low enough to avoid wasted peer broadcasts"),
        ("Bloom auto-scale", "Yes", "Scale with sqrt(N); prevents saturation at high N"),
        ("Peer routing", "Broadcast", "All peers; simpler than key-hash routing"),
        ("Entry size", "8B", "32b key + 32b value; matches RCuckoo for comparison"),
        ("Server RTT", "3 ticks", "Cross-rack model; peer RTT = 1 tick"),
    ]

    tb = sl.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9.0), Inches(5.5))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (c1, c2, c3) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{c1:<18s}  {c2:<12s}  {c3}"
        p.font.size = Pt(13) if i > 0 else Pt(14)
        p.font.name = FONT_MONO
        p.font.bold = (i == 0)
        p.font.color.rgb = TEAL if i == 0 else DARK_GRAY
        p.space_after = Pt(6)
        if i == 0:
            p2 = tf.add_paragraph()
            p2.text = "\u2500" * 70
            p2.font.size = Pt(8)
            p2.font.name = FONT_MONO
            p2.font.color.rgb = LIGHT_GRAY

    # ━━ 15. SECTION: EVALUATION ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    sn += 1
    _section_slide(prs, "Evaluation",
                   "Simulation Framework & Results",
                   CHARCOAL)

    # ━━ 16. SETUP ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    sn += 1
    sl = _content_slide(prs, sn, "Evaluation: Simulation Framework")
    _bullets(sl, [
        "Tick-based discrete event simulation (1 tick = 1 us RDMA RTT)",
        "Identical pre-generated Zipf workloads for both systems",
        "YCSB suite: A (50/50 R/W), B (95/5), C (100% read)",
        "10M entries, 8B per entry, Zipf theta = 0.99",
        "Client counts: {8, 16, 40, 80, 160, 320}",
        "20ms simulated time, 50K ops/client, averaged over 3 trials",
    ], fs=16, sub={
        0: ["Same engine architecture for both systems: fair comparison"],
    })

    _box(sl, Inches(5.5), Inches(4.8), Inches(4.0), Inches(1.5),
         "Latency Model\n\n"
         "Local cache:  0 ticks (free)\n"
         "Peer (same rack):  1 tick\n"
         "Server (cross-rack):  3 ticks",
         WARM_WHITE, 12, False, DARK_GRAY, border=CHARCOAL, shadow=True)

    _note(sl, "Both systems share workload generation, Zipf sampling, and measurement infrastructure.")

    # ━━ 17-23. PLOT SLIDES ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    plot_slides = [
        ("Workload: Zipf Key Distribution", "plot1_zipf_distribution.png", [
            "Zipf(0.99): extreme skew, top-100 keys dominate access",
            "Small LRU cache can serve a large fraction of requests",
            "Same pre-generated workload used by both systems",
        ]),
        ("Throughput (MOPS vs. Clients)", "plot2_throughput_comparison.png", [
            "YCSB-C: OfflineState ~300 vs RCuckoo ~107 MOPS at 320 clients (2.8\u00d7)",
            "YCSB-A: OfflineState ~148 vs RCuckoo ~20 MOPS (7.4\u00d7)",
            "RCuckoo plateaus from lock contention; OfflineState scales near-linearly",
        ]),
        ("Read Source Breakdown", "plot3_read_source_breakdown.png", [
            "~31% local (0 RDMA), ~25% peer (1 RDMA), ~44% server",
            "56% of reads bypass the server entirely",
            "Peer share grows with client count (more collective cached content)",
        ]),
        ("RDMA Calls per Operation", "plot4_rdma_per_op.png", [
            "RCuckoo reads: fixed 1.0 RDMA/op; writes rise to ~2.0",
            "OfflineState YCSB-C: 0.6\u20130.8 RDMA/op (local hits are free)",
            "Fewer RDMA = less network bandwidth, lower tail latency",
        ]),
        ("Read Latency CDF (320 clients)", "plot5_latency_cdf.png", [
            "OfflineState: three-tier staircase (0, 1, 3 ticks)",
            "31% at 0 ticks (local), 56% at \u22641 tick; RCuckoo: all at 3 ticks",
            "Staircase shape validates three-tier cache hierarchy",
        ]),
        ("Bloom Filter False Positive Rate", "plot6_bloom_fpr.png", [
            "Auto-scaling keeps FPR below 5% at 320 clients",
            "Without scaling: FPR > 40% from bloom saturation",
            "Theoretical FPR matches empirical (validates sizing formula)",
        ]),
        ("Sensitivity: Zipf Skewness", "plot7_skewness_sensitivity.png", [
            "OfflineState gains more from higher skewness (better cache hits)",
            "At theta=0.5: systems converge; caching has little benefit",
            "At theta=0.99: OfflineState RDMA/op \u2248 0.6 vs RCuckoo fixed 1.0",
        ]),
    ]

    for title, img, takeaway in plot_slides:
        sn += 1
        sl = _content_slide(prs, sn, title, size=24)
        _image(sl, img)
        _insight_strip(sl, takeaway)

    # ━━ 24. SYNC + CACHE SENSITIVITY ━━━━━━━━━━━━━━━━━━━━━━━━
    sn += 1
    sl = _content_slide(prs, sn, "Sensitivity: Sync Interval & Cache Size", size=24)
    _image(sl, "plot8_sync_sensitivity.png", top=Inches(0.8), height=Inches(2.65))
    _image(sl, "plot9_cache_sensitivity.png", top=Inches(3.55), height=Inches(2.65))
    _insight_strip(sl, [
        "Sync interval: robust across 10\u20131000 ticks; FPR rises steeply beyond",
        "Cache size: 256KB saturates benefit; diminishing returns above",
        "Both parameters have wide sweet spots \u2014 design is not fragile",
    ], top=Inches(6.3), bg=TEAL_LIGHT)

    # ━━ 25. CONCLUSIONS ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    sn += 1
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, CHARCOAL)

    tb = sl.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "Conclusions"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.name = FONT_BODY
    p.font.color.rgb = WHITE

    findings = [
        (STEEL,
         "Same-rack: RCuckoo wins",
         "Single-RTT reads, no bloom dependency. When server \u2248 peer latency, caching adds overhead."),
        (TEAL,
         "Cross-rack: OfflineState wins (up to 2.8\u00d7)",
         "Caching eliminates most server trips. 56% of reads bypass server at 320 clients."),
        (ORANGE,
         "Write-heavy: OfflineState wins (up to 7.4\u00d7)",
         "Lock-free writes avoid contention. RCuckoo MCAS retries grow quadratically."),
    ]

    for i, (color, title, desc) in enumerate(findings):
        y = Inches(1.4) + Inches(i * 1.3)
        bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.8), y, Inches(0.08), Inches(1.0))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        tb = sl.shapes.add_textbox(Inches(1.1), y, Inches(8.1), Inches(1.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.name = FONT_BODY
        p.font.color.rgb = color
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.name = FONT_BODY
        p2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

    # Future work
    ln = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0.8), Inches(5.1), Inches(8.4), Pt(1))
    ln.fill.solid()
    ln.fill.fore_color.rgb = TEAL
    ln.line.fill.background()

    tb = sl.shapes.add_textbox(Inches(0.8), Inches(5.25), Inches(8.4), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.text = "Future Work"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.name = FONT_BODY
    p.font.color.rgb = TEAL

    future = [
        "Hardware validation on multi-rack RDMA testbed",
        "Adaptive bloom sizing based on runtime workload",
        "Hybrid: OfflineState caching + RCuckoo index",
    ]
    tb2 = sl.shapes.add_textbox(Inches(0.8), Inches(5.6), Inches(8.4), Inches(0.9))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    for i, item in enumerate(future):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = "\u2013  " + item
        p.font.size = Pt(12)
        p.font.name = FONT_BODY
        p.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)

    # References
    refs = (
        '[1] Grant & Snoeren, ATC\'25  '
        '[2] Shen et al. (FUSEE), FAST\'23  '
        '[3] Tsai et al. (Clover), ATC\'20  '
        '[4] Wang et al. (Sherman), SIGMOD\'22  '
        '[5] Li et al., EuroSys\'14  '
        '[6] Fan et al. (MemC3), NSDI\'13  '
        '[7] Kalia et al., SIGCOMM\'14  '
        '[8] Pagh & Rodler, J.Alg.\'04  '
        '[9] Bloom, CACM\'70'
    )
    tb3 = sl.shapes.add_textbox(Inches(0.3), Inches(6.7), Inches(9.4), Inches(0.6))
    p = tb3.text_frame.paragraphs[0]
    p.text = refs
    p.font.size = Pt(8)
    p.font.name = FONT_BODY
    p.font.color.rgb = RGBColor(0x77, 0x77, 0x77)

    _footer(sl, sn)

    # ── Save ──────────────────────────────────────────────────
    out = os.path.join(BASEDIR, "presentation.pptx")
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Slides: {sn}")
    return out


if __name__ == "__main__":
    build()
